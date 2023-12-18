import os
from PyPDF2 import PdfReader
from pptx import Presentation

def save_uploaded_file(folder_path: str, file):
    file_path = os.path.join(folder_path, file.filename)
    with open(file_path, "wb") as buffer:
        buffer.write(file.file.read())
    return file_path

def extract_text_from_pdf(pdf_path: str) -> str:
    text = ""
    with open(pdf_path, "rb") as file:
        pdf_reader = PdfReader(file)
        for page_num in range(len(pdf_reader.pages)):
            text += pdf_reader.pages[page_num].extract_text()
    return text

def extract_text_from_ppt(ppt_path: str) -> str:
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
