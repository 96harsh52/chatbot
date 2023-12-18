import os
import difflib
import pandas as pd
from openai import OpenAI
from datetime import datetime
from PyPDF2 import PdfReader
from pptx import Presentation

# Set OpenAI API key

os.environ["OPENAI_API_KEY"] = "Api-key"
# Initialize OpenAI client
client = OpenAI()
=====

def read_text_file(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        return f"Error reading file {file_path}: {str(e)}"

def is_similar(text1, text2, threshold=0.8):
    try:
        d = difflib.SequenceMatcher(None, text1, text2)
        return d.ratio() >= threshold
    except Exception as e:
        return f"Error comparing texts: {str(e)}"

def get_embedding(text, model="text-embedding-ada-002"):
    try:
        text = text.replace("\n", " ")
        res = client.embeddings.create(input=[text], model=model)
        return res.data[0].embedding
    except Exception as e:
        with open(f'logs.log', 'a', encoding='utf-8') as log_file:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_file.write(f"[{timestamp}] Error getting embedding: {str(e)}\n")
        raise ValueError(f"Error getting embedding: {str(e)}")

def save_uploaded_file(folder_path: str, file):
    file_path = os.path.join(folder_path, file.filename)
    with open(file_path, "wb") as buffer:
        buffer.write(file.file.read())
    return file_path

def extract_text_from_pdf(pdf_path: str) -> str:
    text = ""
    with open(pdf_path, "rb") as file:
        pdf_reader = PdfReader(file)
        for page_num in range(len(pdf_reader.pages)):
            text += pdf_reader.pages[page_num].extract_text()
    return text

def extract_text_from_ppt(ppt_path: str) -> str:
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def process_files(input_directory, output_directory):
    try:
        df_data = []
        file_contents = {}
        os.makedirs(output_directory, exist_ok=True)
        folder_name = os.path.basename(input_directory)
        num_files_processed = 0
        num_failed_embeddings = 0

        for filename in os.listdir(input_directory):
            if filename.endswith(('.txt', '.docx', '.doc', '.csv', '.json', '.pdf', '.ppt', '.pptx')):
                file_path = os.path.join(input_directory, filename)

                # Extract text based on file extension
                if filename.lower().endswith('.pdf'):
                    file_text = extract_text_from_pdf(file_path)
                elif filename.lower().endswith(('.ppt', '.pptx')):
                    file_text = extract_text_from_ppt(file_path)
                else:
                    file_text = read_text_file(file_path)

                include_file = True
                for existing_filename, existing_text in file_contents.items():
                    if is_similar(file_text, existing_text):
                        include_file = False
                        break

                if include_file:
                    embedding = get_embedding(file_text, model='text-embedding-ada-002')

                    if embedding is None:
                        with open(f'logs.log', 'a', encoding='utf-8') as log_file:
                            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                            log_file.write(f"[{timestamp}] Skipping file {filename} due to embedding error.\n")
                        num_failed_embeddings += 1
                        continue

                    file_contents[filename] = file_text
                    df_data.append([filename, file_text, embedding])
                    num_files_processed += 1

        if num_files_processed == 0:
            return {"status": 401, "message": f"Error: No valid embeddings obtained. No files processed in {folder_name}."}

        columns = ['File_Name', 'text', 'embedding']
        df = pd.DataFrame(df_data, columns=columns)
        output_file_path = os.path.join(output_directory, f'{folder_name}_embedded.csv')
        df.to_csv(output_file_path, index=False)

        with open(f'logs.log', 'a', encoding='utf-8') as log_file:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_file.write(f"[{timestamp}] Successfully processed {num_files_processed} files in {folder_name}. "
                           f"CSV file '{output_file_path}' created.\n")
            if num_failed_embeddings > 0:
                log_file.write(f"[{timestamp}] {num_failed_embeddings} files in {folder_name} had embedding creation failures.\n")

        return {"status": 200, "message": f"CSV file '{output_file_path}' created successfully."}

    except Exception as e:
        with open(f'logs.log', 'a', encoding='utf-8') as log_file:
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            log_file.write(f"[{timestamp}] Error processing files in {folder_name}: {str(e)}\n")
        raise ValueError(f"Error processing files in {folder_name}: {str(e)}")


# import os
# import difflib
# import pandas as pd
# from openai import OpenAI
# from datetime import datetime
#
# # Set OpenAI API key
# os.environ["OPENAI_API_KEY"] = "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"
# # Initialize OpenAI client
# client = OpenAI()
#
# def read_text_file(file_path):
#     try:
#         with open(file_path, 'r', encoding='utf-8') as file:
#             return file.read()
#     except Exception as e:
#         return f"Error reading file {file_path}: {str(e)}"
#
# def is_similar(text1, text2, threshold=0.8):
#     try:
#         d = difflib.SequenceMatcher(None, text1, text2)
#         return d.ratio() >= threshold
#     except Exception as e:
#         return f"Error comparing texts: {str(e)}"
#
# def get_embedding(text, model="text-embedding-ada-002"):
#     try:
#         text = text.replace("\n", " ")
#         res = client.embeddings.create(input=[text], model=model)
#         return res.data[0].embedding
#     except Exception as e:
#         with open(f'logs.log', 'a', encoding='utf-8') as log_file:
#             timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
#             log_file.write(f"[{timestamp}] Error getting embedding: {str(e)}\n")
#         raise ValueError(f"Error getting embedding: {str(e)}")
#
# def process_files(input_directory, output_directory):
#     try:
#         df_data = []
#         file_contents = {}
#         os.makedirs(output_directory, exist_ok=True)
#         folder_name = os.path.basename(input_directory)
#         num_files_processed = 0
#         num_failed_embeddings = 0
#
#         for filename in os.listdir(input_directory):
#             if filename.endswith(".txt"):
#                 file_path = os.path.join(input_directory, filename)
#                 file_text = read_text_file(file_path)
#
#                 include_file = True
#                 for existing_filename, existing_text in file_contents.items():
#                     if is_similar(file_text, existing_text):
#                         include_file = False
#                         break
#
#                 if include_file:
#                     embedding = get_embedding(file_text, model='text-embedding-ada-002')
#
#                     if embedding is None:
#                         with open(f'logs.log', 'a', encoding='utf-8') as log_file:
#                             timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
#                             log_file.write(f"[{timestamp}] Skipping file {filename} due to embedding error.\n")
#                         num_failed_embeddings += 1
#                         continue
#
#                     file_contents[filename] = file_text
#                     df_data.append([filename, file_text, embedding])
#                     num_files_processed += 1
#
#         if num_files_processed == 0:
#             return {"status": 401, "message": f"Error: No valid embeddings obtained. No files processed in {folder_name}."}
#
#         columns = ['File_Name', 'text', 'embedding']
#         df = pd.DataFrame(df_data, columns=columns)
#         output_file_path = os.path.join(output_directory, f'{folder_name}_embedded.csv')
#         df.to_csv(output_file_path, index=False)
#
#         with open(f'logs.log', 'a', encoding='utf-8') as log_file:
#             timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
#             log_file.write(f"[{timestamp}] Successfully processed {num_files_processed} files in {folder_name}. "
#                            f"CSV file '{output_file_path}' created.\n")
#             if num_failed_embeddings > 0:
#                 log_file.write(f"[{timestamp}] {num_failed_embeddings} files in {folder_name} had embedding creation failures.\n")
#
#         return {"status": 200, "message": f"CSV file '{output_file_path}' created successfully."}
#
#     except Exception as e:
#         with open(f'logs.log', 'a', encoding='utf-8') as log_file:
#             timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
#             log_file.write(f"[{timestamp}] Error processing files in {folder_name}: {str(e)}\n")
#         raise ValueError(f"Error processing files in {folder_name}: {str(e)}")



# Example usage
# if __name__ == "__main__":
#     input_directory = r'D:\CG-Vks\chatbot_10\data\www.allamericanpha.com'
#     output_directory = r'D:\CG-Vks\chatbot_10\New folder'
#     print(process_files(input_directory, output_directory))
